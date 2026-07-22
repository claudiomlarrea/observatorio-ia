#!/usr/bin/env python3
"""Plantilla PPTX de 6 diapositivas editables para Jornadas IA 2026."""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_REPO = ROOT / "docs/plantillas/plantilla-presentacion-jornadas-ia-2026.pptx"
OUT_ONEDRIVE = Path(
    "/Users/claudiolarrea/Library/CloudStorage/OneDrive-Personal/16 Secretaría de Investigación/"
    "60 Observatorio de Inteligencia Artificial/Jornadas de IA 2026/"
    "plantilla-presentacion-jornadas-ia-2026.pptx"
)
LOGO = ROOT / "assets/logo-observatorio-ia.png"
LOGO_CIRCLE = ROOT / "assets/logo-observatorio-ia-circle.png"

GREEN = RGBColor(0x04, 0x2F, 0x23)
GREEN_MID = RGBColor(0x06, 0x4A, 0x38)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x1F, 0x14, 0x18)
MUTED = RGBColor(0x5C, 0x4F, 0x54)
SAGE = RGBColor(0xF4, 0xF8, 0xF6)


def _ensure_circle_logo() -> Path:
    if LOGO_CIRCLE.is_file():
        return LOGO_CIRCLE
    from PIL import Image, ImageDraw

    im = Image.open(LOGO).convert("RGBA")
    size = min(im.size)
    im = im.resize((size, size), Image.Resampling.LANCZOS)
    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, size - 1, size - 1), fill=255)
    out = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    out.paste(im, (0, 0))
    out.putalpha(mask)
    LOGO_CIRCLE.parent.mkdir(parents=True, exist_ok=True)
    out.save(LOGO_CIRCLE, "PNG")
    return LOGO_CIRCLE


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_line(shape) -> None:
    shape.line.fill.background()


def _write_lines(shape, lines):
    """lines: list of (text, size, bold, color[, align])."""
    tf = shape.text_frame
    tf.word_wrap = True
    # vaciar párrafos existentes
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first_p = tf.paragraphs[0]
    for r in list(first_p.runs):
        r._r.getparent().remove(r._r)
    first_p.text = ""

    for i, item in enumerate(lines):
        text, size, bold, color = item[:4]
        align = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        p = first_p if i == 0 else tf.add_paragraph()
        if i == 0:
            p.text = ""
        p.alignment = align
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _chrome(slide, prs, logo_path: Path, page: int, total: int = 6) -> None:
    # fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _fill(bg, WHITE)

    banner_h = Inches(1.15)
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, banner_h)
    _fill(banner, GREEN)

    gold = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, banner_h, prs.slide_width, Inches(0.05)
    )
    _fill(gold, GOLD)

    logo_size = Inches(0.9)
    slide.shapes.add_picture(
        str(logo_path), Inches(0.3), Inches(0.12), logo_size, logo_size
    )

    title_box = slide.shapes.add_textbox(Inches(1.4), Inches(0.18), Inches(10.5), Inches(0.9))
    _write_lines(
        title_box,
        [
            ("UNIVERSIDAD CATÓLICA DE CUYO · Observatorio de IA", 11, False, WHITE),
            ("1° Jornadas internas de Inteligencia Artificial 2026", 18, True, WHITE),
            ("Encuentro virtual · 6 de octubre de 2026 · 15:00 h", 12, False, WHITE),
        ],
    )

    footer = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        prs.slide_height - Inches(0.4),
        prs.slide_width,
        Inches(0.4),
    )
    _fill(footer, GREEN_MID)
    foot_tb = slide.shapes.add_textbox(
        Inches(0.4),
        prs.slide_height - Inches(0.36),
        Inches(11.5),
        Inches(0.32),
    )
    _write_lines(
        foot_tb,
        [
            (
                f"Observatorio de IA · UCCuyo  ·  Diapositiva {page}/{total}  ·  "
                "https://claudiomlarrea.github.io/observatorio-ia/#jornadas-ia",
                11,
                False,
                WHITE,
                PP_ALIGN.CENTER,
            )
        ],
    )


def _section_title(slide, text: str):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(0.55))
    _write_lines(box, [(text, 26, True, GREEN)])


def _body_card(slide, lines):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(2.1),
        Inches(12.1),
        Inches(4.6),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = SAGE
    card.line.color.rgb = GREEN_MID
    card.line.width = Pt(1.25)
    _write_lines(card, lines)


def build() -> Path:
    logo_path = _ensure_circle_logo()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 — Portada
    s1 = prs.slides.add_slide(blank)
    _chrome(s1, prs, logo_path, 1)
    _section_title(s1, "1. Portada")
    _body_card(
        s1,
        [
            ("Título de la ponencia", 28, True, GREEN),
            ("(hacé clic aquí y reemplazá este texto)", 14, False, MUTED),
            ("", 10, False, MUTED),
            ("Autores: N. Apellido¹, N. Apellido²", 18, True, TEXT),
            ("Unidad académica / Facultad o Instituto", 16, False, MUTED),
            ("correo@uccuyo.edu.ar", 14, False, MUTED),
            ("", 10, False, MUTED),
            (
                "Cómo usar esta plantilla: hay 6 diapositivas editables. "
                "Hacé clic en cada texto para modificarlo. Duplicá diapositivas solo si hace falta "
                "(máximo 6 en total).",
                13,
                False,
                MUTED,
            ),
        ],
    )

    # 2 — Objetivo / contexto
    s2 = prs.slides.add_slide(blank)
    _chrome(s2, prs, logo_path, 2)
    _section_title(s2, "2. Objetivo / contexto")
    _body_card(
        s2,
        [
            ("¿Qué problema o pregunta aborda el trabajo?", 18, True, GREEN),
            ("• Escribí aquí el objetivo general", 16, False, TEXT),
            ("• Contexto institucional o del área", 16, False, TEXT),
            ("• Por qué importa para la UCCuyo / la región", 16, False, TEXT),
        ],
    )

    # 3 — Método
    s3 = prs.slides.add_slide(blank)
    _chrome(s3, prs, logo_path, 3)
    _section_title(s3, "3. Método o experiencia")
    _body_card(
        s3,
        [
            ("¿Cómo se hizo?", 18, True, GREEN),
            ("• Enfoque o diseño (docencia, investigación, gestión…)", 16, False, TEXT),
            ("• Herramientas / datos / participantes", 16, False, TEXT),
            ("• Pasos principales de la experiencia", 16, False, TEXT),
        ],
    )

    # 4 — Resultados
    s4 = prs.slides.add_slide(blank)
    _chrome(s4, prs, logo_path, 4)
    _section_title(s4, "4. Resultados / hallazgos")
    _body_card(
        s4,
        [
            ("¿Qué se obtuvo?", 18, True, GREEN),
            ("• Hallazgo principal 1", 16, False, TEXT),
            ("• Hallazgo principal 2", 16, False, TEXT),
            ("• Dato o ejemplo concreto (sin saturar la diapositiva)", 16, False, TEXT),
        ],
    )

    # 5 — Conclusiones
    s5 = prs.slides.add_slide(blank)
    _chrome(s5, prs, logo_path, 5)
    _section_title(s5, "5. Conclusiones")
    _body_card(
        s5,
        [
            ("Mensajes clave para llevarse", 18, True, GREEN),
            ("• Conclusión 1", 16, False, TEXT),
            ("• Conclusión 2", 16, False, TEXT),
            ("• Implicancias o próximos pasos", 16, False, TEXT),
        ],
    )

    # 6 — Cierre
    s6 = prs.slides.add_slide(blank)
    _chrome(s6, prs, logo_path, 6)
    _section_title(s6, "6. Cierre")
    _body_card(
        s6,
        [
            ("Gracias", 32, True, GREEN),
            ("", 10, False, MUTED),
            ("Preguntas y comentarios", 18, False, TEXT),
            ("correo@uccuyo.edu.ar", 16, False, MUTED),
            ("", 10, False, MUTED),
            (
                "Exposición: hasta 10 minutos · Máximo 6 diapositivas · "
                "Carga hasta el 10 de septiembre de 2026",
                14,
                False,
                MUTED,
            ),
        ],
    )

    OUT_REPO.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_REPO))
    return OUT_REPO


def main() -> None:
    path = build()
    print(f"Wrote {path} ({path.stat().st_size} bytes)")
    prs = Presentation(str(path))
    print(f"Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides, 1):
        types = [str(s.shape_type).split(".")[-1] for s in slide.shapes]
        print(f"  {i}: {len(slide.shapes)} shapes → {', '.join(types[:8])}...")
    if OUT_ONEDRIVE.parent.is_dir():
        shutil.copy2(path, OUT_ONEDRIVE)
        print(f"Copied to {OUT_ONEDRIVE}")


if __name__ == "__main__":
    main()
